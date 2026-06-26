from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # -----------------------------
    # Helper to add a content slide
    # -----------------------------
    def add_slide(title_text, bullet_points):
        # 1 is the layout for Title and Content
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = title_text
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        for idx, point in enumerate(bullet_points):
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            # Basic formatting
            p.font.size = Pt(18)
            # Add spacing between paragraphs
            p.space_after = Pt(12)

    # ==========================================
    # Slide 1: Title & Problem Statement
    # ==========================================
    # 0 is the layout for Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Talent OS: AI-Powered Talent Intelligence"
    subtitle.text = (
        "Team Name: Neural Nexus (or Your Team Name)\n"
        "Team Leader Name: Priyam\n\n"
        "Problem Statement: ATS systems act as static repositories, forcing recruiters to manually parse "
        "unstructured data. This leads to missed 'hidden talent', bias, and an inability to objectively "
        "benchmark candidates against dynamic JDs and market standards."
    )

    # ==========================================
    # Slide 2: Solution Overview
    # ==========================================
    add_slide(
        "Solution Overview",
        [
            "Proposed Solution: 'Talent OS' - an enterprise-grade AI Talent Intelligence Platform built with Python and Streamlit.",
            "It actively scores, segments, and benchmarks candidates via an interactive dashboard.",
            "Differentiation: Instead of standard keyword search, Talent OS uses a multi-faceted evaluation engine.",
            "It computes 'Learning Velocity', 'Market Value Prediction', and 'Hidden Talent' synergy via semantic embedding models."
        ]
    )

    # ==========================================
    # Slide 3: JD Understanding & Candidate Evaluation
    # ==========================================
    add_slide(
        "JD Understanding & Candidate Evaluation",
        [
            "Key Requirements Extracted: Core technical skills, required experience level, and nuanced contextual requirements from the JD.",
            "Key Candidate Signals: Career trajectory (duration, roles), skill proficiency (endorsements), education quality, and cognitive assessment scores.",
            "Evaluation Beyond Keywords: We use Sentence Transformers to map the semantic intent of the JD against the candidate's holistic profile (Summary + Headline).",
            "This ensures candidates who use contextual synonyms or have transferable skills aren't overlooked."
        ]
    )

    # ==========================================
    # Slide 4: Ranking Methodology
    # ==========================================
    add_slide(
        "Ranking Methodology",
        [
            "System Retrieval & Scoring: Uses a composite 'Talent Score' combining semantic match, learning velocity, and market benchmarking.",
            "Models & Algorithms: Leveraging pre-trained NLP (HuggingFace Sentence-Transformers) for dense embedding comparisons.",
            "Heuristics: Rule-based weights for skill proficiency arrays and career gap penalties.",
            "Signal Combination: The final ranking is a weighted sum: 40% Semantic Match, 30% Learning Velocity, 30% Expected vs Market Salary."
        ]
    )

    # ==========================================
    # Slide 5: Explainability & Data Validation
    # ==========================================
    add_slide(
        "Explainability & Data Validation",
        [
            "Explainability: The UI provides a 'Score Breakdown' Donut chart and explicitly lists 'AI Match Reasoning' (High/Moderate/Low Semantic Overlap).",
            "Preventing Hallucinations: We do not use generative LLMs to make up facts; scoring is strictly bounded by cosine similarity and deterministic math.",
            "Handling Low-Quality Profiles: The Data Cleaning pipeline (data_cleaning.py) aggressively standardizes inputs, drops invalid signals, and penalizes candidates with empty fields via the 'profile_completeness' metric."
        ]
    )

    # ==========================================
    # Slide 6: End-to-End Workflow
    # ==========================================
    add_slide(
        "End-to-End Workflow",
        [
            "1. Ingestion: 'candidates_df.py' parses nested JSON candidate data.",
            "2. Standardization: 'data_cleaning.py' normalizes text and handles missing values.",
            "3. Engineering: 'feature_engineering.py' calculates compound metrics (e.g., tenure, skill density).",
            "4. ML Scoring: User pastes a JD into the Streamlit UI.",
            "5. Ranking: Candidate embeddings are compared to JD embedding.",
            "6. Visualization: Output displayed on the interactive Enterprise Dashboard."
        ]
    )

    # ==========================================
    # Slide 7: System Architecture
    # ==========================================
    add_slide(
        "System Architecture",
        [
            "Frontend: Streamlit (with heavily customized enterprise CSS & Bento Box layouts).",
            "Backend/Data Processing: Pandas & NumPy for relational data merging and high-speed vectorized operations.",
            "Machine Learning Layer: Scikit-learn (similarity metrics) & Sentence-Transformers (NLP).",
            "Data Storage: Local JSON/CSV pipeline designed for scalability to cloud storage (S3/GCS)."
        ]
    )

    # ==========================================
    # Slide 8: Results & Performance
    # ==========================================
    add_slide(
        "Results & Performance",
        [
            "Ranking Quality Insights: The semantic engine successfully identifies candidates with 'latent' synergies even if they lack exact JD keywords.",
            "Runtime Constraints: By using lightweight embedding models and vectorized Pandas operations, the system ranks thousands of candidates in sub-second times.",
            "UI Performance: Custom HTML/CSS injections maintain 60FPS rendering in the browser without requiring a heavy React/Next.js frontend."
        ]
    )

    # ==========================================
    # Slide 9: Technologies Used
    # ==========================================
    add_slide(
        "Technologies Used",
        [
            "Python 3: The core language for unified data science and web deployment.",
            "Pandas & NumPy: Selected for robust, vectorized tabular data manipulation.",
            "Streamlit: Selected for rapid prototyping of analytical web apps with pure Python.",
            "Plotly: Selected for highly interactive, enterprise-grade data visualizations (Gauge & Donut charts).",
            "Hugging Face (Sentence-Transformers): Selected for state-of-the-art semantic matching without the latency/cost of external API calls (e.g., OpenAI)."
        ]
    )

    # ==========================================
    # Slide 10: Submission Assets
    # ==========================================
    add_slide(
        "Submission Assets",
        [
            "GitHub Repository: [Insert Link to your repo]",
            "Video Demo: [Insert Link to Demo Video]",
            "Live App URL (Optional): [Insert Streamlit Cloud/Render URL]",
            "Contact: [Insert your email]"
        ]
    )

    prs.save('Talent_OS_Presentation.pptx')
    print("Presentation saved as Talent_OS_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()
